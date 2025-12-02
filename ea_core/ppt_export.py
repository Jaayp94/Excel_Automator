# ea_core/ppt_export.py

import logging
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

logger = logging.getLogger(__name__)


def export_analysis_to_pptx(filepath: str, df: pd.DataFrame, title: str = "", subtitle: str = "") -> Path:
    """
    Erstellt einen PowerPoint-Report:
    1. Titelfolie
    2. Balkendiagramm
    3. Tortendiagramm
    4. Tabelle
    """
    if df is None or df.empty:
        raise ValueError("Analysedaten leer!")

    cols = list(df.columns)
    if len(cols) < 2:
        raise ValueError("Mindestens 2 Spalten für Diagramme notwendig.")

    cat_col = cols[0]
    val_col = cols[1]

    path = Path(filepath)
    path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    # ---- Titelfolie ----
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title or "Datenreport"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle or "Automatisch erstellt"

    max_points = 20
    cats = df[cat_col].astype(str).tolist()[:max_points]
    vals = pd.to_numeric(df[val_col], errors="coerce").fillna(0).tolist()[:max_points]

    # ---- Balkendiagramm ----
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Balkendiagramm"

    chart_data = CategoryChartData()
    chart_data.categories = cats
    chart_data.add_series(str(val_col), vals)

    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.5), Inches(9), Inches(4.5),
        chart_data
    )

    # ---- Tortendiagramm ----
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Tortendiagramm"

    pie_data = ChartData()
    pie_data.categories = cats
    pie_data.add_series(str(val_col), vals)

    slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(1), Inches(1.5), Inches(8), Inches(4.5),
        pie_data
    )

    # ---- Tabelle ----
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Tabellarische Übersicht"

    preview = df.head(10)
    rows = len(preview) + 1
    cols_n = len(cols)

    table_shape = slide.shapes.add_table(
        rows, cols_n,
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(4.5),
    )
    table = table_shape.table

    for j, c in enumerate(cols):
        table.cell(0, j).text = str(c)

    for i in range(len(preview)):
        for j, c in enumerate(cols):
            v = preview.iloc[i, j]
            table.cell(i + 1, j).text = "" if pd.isna(v) else str(v)

    prs.save(path)
    return path
